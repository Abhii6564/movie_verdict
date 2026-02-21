from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- STYLE CONSTANTS ---
COLOR_BG = RGBColor(255, 255, 255)      # White
COLOR_PURPLE = RGBColor(102, 51, 153)   # Rebecca Purple
COLOR_TEXT = RGBColor(40, 40, 40)       # Dark Grey/Black
COLOR_WHITE = RGBColor(255, 255, 255)   # White
COLOR_ACCENT = RGBColor(230, 210, 250)    # Light Lavender for backgrounds

def apply_base_design(slide, title_text):
    """Applies a professional light theme with purple accents."""
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    # Header Rectangle (Purple)
    header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.0))
    header_rect.fill.solid()
    header_rect.fill.fore_color.rgb = COLOR_PURPLE
    header_rect.line.fill.background()
    
    # Title Text
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title_text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            
def add_designed_slide(prs, title_text, bullets=None, image_path=None, layout_idx=1):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    apply_base_design(slide, title_text)
    
    # Adjust Content Placeholder
    if bullets:
        placeholder = slide.placeholders[1]
        placeholder.top = Inches(1.4)
        placeholder.left = Inches(0.5)
        placeholder.width = Inches(5) if image_path else Inches(9)
        
        tf = placeholder.text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.color.rgb = COLOR_TEXT
            p.font.size = Pt(22)
            p.level = 0
            p.space_after = Pt(15)
            
    # Add Image on the right if provided
    if image_path and os.path.exists(image_path):
        left = Inches(5.8)
        top = Inches(1.8)
        width = Inches(3.8)
        slide.shapes.add_picture(image_path, left, top, width=width)

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0), Inches(10), Inches(3.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_PURPLE
    rect.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "MOVIE SUCCESS PREDICTION"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = COLOR_WHITE
    p.font.size = Pt(54)
    p.font.bold = True
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    stf = sub_box.text_frame
    stf.text = "Data-Driven Insights & Financial Forecasting"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.color.rgb = COLOR_WHITE
    sp.font.size = Pt(28)
    sp.font.italic = True

    by_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    btf = by_box.text_frame
    btf.text = "Project Presentation | Predictive Analytics for Cinema"
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.color.rgb = COLOR_TEXT
    bp.font.size = Pt(20)
    bp.font.bold = True

    # Slide 2: Executive Summary
    add_designed_slide(prs, "Executive Summary", [
        "Objective: Quantify the 'hit' potential of film projects using historical data.",
        "Target: Triple-Class ROI Status (Hit, Average, Flop).",
        "Key Differentiator: Financial outcome prediction vs. standard rating prediction.",
        "Value: Reducing investment risk in the volatile global film market."
    ])

    # Slide 3: Dataset Deep Dive
    add_designed_slide(prs, "The Dataset: 5,000+ Movies", [
        "Source: IMDB Metadata consisting of 28 core attributes.",
        "Dimensions: Budget, Gross Revenue, Social Media Metrics, Genre, Cast.",
        "Data Quality: Rigorous cleaning removed 100+ duplicates and handled missing 10-20% values.",
        "Target Logic: Calculated ROI Ratio = Gross / Budget."
    ])

    # Slide 4: Data Distribution & Skewness
    add_designed_slide(prs, "Historical Distribution & Imbalance", [
        "Data Skew: A large majority of recorded movies are 'Hits' in this dataset.",
        "Budget Variance: Ranging from independent films to $300M blockbusters.",
        "Revenue Gaps: Significant variance between theatrical releases and streaming success.",
        "Goal: Training the model to detect 'Flops' despite sample sparsity."
    ], "chart_distribution.png")

    # Slide 5: Feature Engineering Architecture
    add_designed_slide(prs, "The Success Drivers", [
        "Financial Weight: Budget proves to be the strongest success predictor.",
        "Social Influence: Cast FB Likes represent 'star power' and initial hype.",
        "Genre Influence: Critical for determining initial market size and budget caps.",
        "Scaling: Data normalized via StandardScaler for neural model compatibility."
    ], "chart_importance.png")

    # Slide 6: Market Performance Trends
    add_designed_slide(prs, "Market Trends: The 20-Year Shift", [
        "Hyper-Inflation: Movie budgets have scaled significantly over the last two decades.",
        "Revenue Consolidation: Major hits are earning more, while averages are shrinking.",
        "Risk Factor: Modern movies require higher Gross to reach 'Hit' status (2x ROI barrier)."
    ], "chart_trends.png")

    # Slide 7: Predictive Results
    add_designed_slide(prs, "Model Performance (Random Forest)", [
        "Algorithm: Random Forest with 100 decision trees.",
        "Optimization: Used Class Weights to penalize misclassification of rare flops.",
        "Hit Accuracy: Extremely reliable in predicting successful movie outcomes.",
        "Future Refinement: Moving towards SMOTE for synthetic flop data generation."
    ])

    # Slide 8: Interactive Application
    add_designed_slide(prs, "Interactive Streamlit Deployment", [
        "Real-time 'What-If' analysis for producers and investors.",
        "Feature Input: Genre selection, Budgeting, and Cast Popularity sliders.",
        "Visual Output: Direct probability breakdown and historical context comparison.",
        "Accessibility: Lightweight web-based deployment for executive review."
    ])

    # Slide 9: Future Scope & Roadmap
    add_designed_slide(prs, "The Road Ahead", [
        "Trailer Sentiment Analysis: Extracting hype metrics from social media.",
        "Cast Chemistry: Analyzing co-star compatibility impact on box office.",
        "Advanced Models: transitioning to XGBoost and Deep Neural Networks.",
        "Real-time Integration: Live API feed from Box Office Mojo/Mojo."
    ])

    # Slide 10: Q&A / Close
    add_designed_slide(prs, "THANK YOU", [
        "Project: Movie Success Prediction Analysis",
        "Key Deliverables: Predictive Model & Interactive App",
        "Objective: Data-Driven Support for Film Production Decisions",
        "Questions & Discussion"
    ])

    output_path = "Movie_Success_Prediction_Final.pptx"
    prs.save(output_path)
    print(f"Final Polished Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
